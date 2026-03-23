"""
医局説明会用 PowerPoint スライド作成スクリプト
ANT-DBS マーケティング資材：VNS との有効性比較
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ===== カラーパレット =====
NAVY      = RGBColor(0x00, 0x30, 0x87)   # #003087
ELEC_BLUE = RGBColor(0x00, 0xA8, 0xE0)   # #00A8E0
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF0, 0xF6, 0xFF)   # 薄いブルー背景
ACCENT    = RGBColor(0x00, 0xA8, 0xE0)   # コールアウト用

W = Inches(13.33)   # ワイドスクリーン 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK_LAYOUT = prs.slide_layouts[6]   # 完全空白

# ===== ヘルパー関数 =====

def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width or Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, x, y, w, h,
                font_name="Noto Sans JP", font_size=18,
                bold=False, color=NAVY,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.bold = bold
    run.font.italic = italic
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    return txBox

def add_slide_number(slide, num, total):
    add_textbox(slide, f"{num} / {total}",
                Inches(12.3), Inches(7.1), Inches(1), Inches(0.35),
                font_size=9, color=RGBColor(0xAA, 0xAA, 0xAA),
                align=PP_ALIGN.RIGHT)

def navy_header_bar(slide, title_text):
    """全スライド共通：上部ネイビーバー + タイトル"""
    add_rect(slide, 0, 0, W, Inches(1.1), NAVY)
    add_textbox(slide, title_text,
                Inches(0.4), Inches(0.15), Inches(12.0), Inches(0.85),
                font_size=24, bold=True, color=WHITE)

def elec_line(slide):
    """ヘッダー下部のエレクトリックブルーライン"""
    add_rect(slide, 0, Inches(1.1), W, Inches(0.07), ELEC_BLUE)

def footer_bar(slide):
    """下部フッターライン"""
    add_rect(slide, 0, Inches(7.15), W, Inches(0.04), ELEC_BLUE)
    add_textbox(slide,
                "Zhu J, et al. J Clin Neurosci. 2021;90:112–117  |  本資材は医療従事者向けの情報提供を目的として作成されています",
                Inches(0.3), Inches(7.2), Inches(12.0), Inches(0.25),
                font_size=7, color=RGBColor(0x88, 0x88, 0x88))

TOTAL_SLIDES = 8

# ===========================================================
# スライド 1：表紙
# ===========================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)

# 背景グラデーション風（ネイビーブロック）
add_rect(slide, 0, 0, W, H, NAVY)
add_rect(slide, 0, Inches(4.8), W, Inches(2.7), RGBColor(0x00, 0x1A, 0x4A))

# ブランドアクセントライン
add_rect(slide, Inches(0.4), Inches(1.5), Inches(0.12), Inches(2.2), ELEC_BLUE)

# メインコピー
add_textbox(slide,
            "薬剤抵抗性てんかんに、\nより確かな選択を。",
            Inches(0.7), Inches(1.4), Inches(10.0), Inches(2.5),
            font_size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# サブタイトル
add_textbox(slide,
            "視床前核DBS（ANT-DBS）と迷走神経刺激療法（VNS）の有効性比較",
            Inches(0.7), Inches(3.9), Inches(10.0), Inches(0.7),
            font_size=18, color=ELEC_BLUE)

# エビデンス出典
add_textbox(slide,
            "Zhu J, et al. Journal of Clinical Neuroscience 90 (2021) 112–117",
            Inches(0.7), Inches(4.55), Inches(9.0), Inches(0.4),
            font_size=11, color=RGBColor(0xCC, 0xCC, 0xCC), italic=True)

# タグライン
add_textbox(slide, "Engineering the extraordinary",
            Inches(0.7), Inches(6.7), Inches(6.0), Inches(0.4),
            font_size=13, color=RGBColor(0xAA, 0xCC, 0xFF), italic=True)

# Medtronic ロゴテキスト
add_textbox(slide, "Medtronic",
            Inches(10.5), Inches(6.7), Inches(2.5), Inches(0.4),
            font_size=20, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

# ===========================================================
# スライド 2：疾患背景
# ===========================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
navy_header_bar(slide, "疾患背景：今もコントロールできていない患者がいる")
elec_line(slide)
footer_bar(slide)
add_slide_number(slide, 2, TOTAL_SLIDES)

# 大きな数字コールアウト
add_rect(slide, Inches(0.4), Inches(1.35), Inches(3.8), Inches(2.5), ELEC_BLUE)
add_textbox(slide, "約 30%",
            Inches(0.4), Inches(1.5), Inches(3.8), Inches(1.3),
            font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "てんかん患者のうち\n薬物療法で発作を\nコントロールできない",
            Inches(0.4), Inches(2.7), Inches(3.8), Inches(0.9),
            font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

# 本文説明
add_textbox(slide,
            "世界に約 7,000万人 のてんかん患者が存在し、そのうち約 3割 が\n薬剤抵抗性てんかん（DRE：Drug-resistant epilepsy）に分類されます。\n\n"
            "2剤以上の抗てんかん薬（AED）で適切に治療しても発作が持続し、\n切除術の適応外または術後も発作が続くケースも少なくありません。",
            Inches(4.5), Inches(1.4), Inches(8.4), Inches(2.1),
            font_size=15, color=NAVY)

# 神経刺激療法の位置づけ
add_rect(slide, Inches(0.4), Inches(4.05), Inches(12.4), Inches(0.06), ELEC_BLUE)
add_textbox(slide, "神経刺激療法は、薬物療法・切除術に次ぐ有力な選択肢です",
            Inches(0.4), Inches(4.2), Inches(12.4), Inches(0.5),
            font_size=17, bold=True, color=NAVY)

# 2治療法比較ボックス
add_rect(slide, Inches(0.4), Inches(4.85), Inches(5.8), Inches(1.7), LIGHT_BG)
add_textbox(slide, "VNS（迷走神経刺激療法）",
            Inches(0.6), Inches(4.95), Inches(5.4), Inches(0.4),
            font_size=14, bold=True, color=NAVY)
add_textbox(slide, "FDA承認：1997年（4歳以上）\n慢性かつ広く普及した選択肢",
            Inches(0.6), Inches(5.35), Inches(5.4), Inches(0.85),
            font_size=13, color=NAVY)

add_rect(slide, Inches(6.5), Inches(4.85), Inches(6.3), Inches(1.7), NAVY)
add_textbox(slide, "ANT-DBS（視床前核深部脳刺激療法）",
            Inches(6.7), Inches(4.95), Inches(5.9), Inches(0.4),
            font_size=14, bold=True, color=ELEC_BLUE)
add_textbox(slide, "FDA承認：2018年（18歳以上）\nより標的化された新しい選択肢",
            Inches(6.7), Inches(5.35), Inches(5.9), Inches(0.85),
            font_size=13, color=WHITE)

# ===========================================================
# スライド 3：試験概要
# ===========================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
navy_header_bar(slide, "試験概要：同一施設・同一基準による厳格な比較")
elec_line(slide)
footer_bar(slide)
add_slide_number(slide, 3, TOTAL_SLIDES)

# 試験デザイン情報
info_rows = [
    ("試験デザイン", "後ろ向き観察研究・単施設"),
    ("施設",        "宣武医院（北京首都医科大学 神経外科・機能神経外科センター）"),
    ("登録期間",    "2013年6月〜2018年7月"),
    ("観察期間",    "術前ベースライン〜術後 12ヵ月"),
    ("評価間隔",    "術後 3・6・9・12ヵ月（外来問診）"),
]
for i, (label, val) in enumerate(info_rows):
    y = Inches(1.35) + i * Inches(0.72)
    add_rect(slide, Inches(0.4), y, Inches(2.8), Inches(0.6), NAVY)
    add_textbox(slide, label, Inches(0.4), y + Inches(0.08), Inches(2.8), Inches(0.5),
                font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(3.25), y, Inches(9.5), Inches(0.6), LIGHT_BG)
    add_textbox(slide, val, Inches(3.35), y + Inches(0.08), Inches(9.3), Inches(0.5),
                font_size=13, color=NAVY)

# 患者コールアウト
y_box = Inches(5.0)
add_rect(slide, Inches(0.4), y_box, Inches(5.8), Inches(1.85), RGBColor(0xE8, 0xF4, 0xFB))
add_textbox(slide, "VNS群", Inches(0.5), y_box + Inches(0.1), Inches(5.6), Inches(0.4),
            font_size=15, bold=True, color=NAVY)
add_textbox(slide,
            "17例\n平均年齢 20.24 ± 11.40歳（5〜41歳）",
            Inches(0.5), y_box + Inches(0.5), Inches(5.6), Inches(1.1),
            font_size=13, color=NAVY)

add_rect(slide, Inches(6.7), y_box, Inches(6.2), Inches(1.85), NAVY)
add_textbox(slide, "ANT-DBS群", Inches(6.8), y_box + Inches(0.1), Inches(6.0), Inches(0.4),
            font_size=15, bold=True, color=ELEC_BLUE)
add_textbox(slide,
            "18例\n平均年齢 28.94 ± 12.00歳（12〜52歳）",
            Inches(6.8), y_box + Inches(0.5), Inches(6.0), Inches(1.1),
            font_size=13, color=WHITE)

add_textbox(slide,
            "全例：長時間ビデオ脳波・MRI・PET・MEG による厳格な術前評価で切除術を除外",
            Inches(0.4), Inches(6.95), Inches(12.4), Inches(0.3),
            font_size=10, color=RGBColor(0x55, 0x55, 0x55), italic=True)

# ===========================================================
# スライド 4：主要結果（発作減少率）
# ===========================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
navy_header_bar(slide, "主要結果：ANT-DBSは全時点でVNSを上回る発作減少率")
elec_line(slide)
footer_bar(slide)
add_slide_number(slide, 4, TOTAL_SLIDES)

# 大数値コールアウト
add_rect(slide, Inches(0.4), Inches(1.35), Inches(3.5), Inches(2.6), NAVY)
add_textbox(slide, "65.28%",
            Inches(0.4), Inches(1.5), Inches(3.5), Inches(1.3),
            font_size=54, bold=True, color=ELEC_BLUE, align=PP_ALIGN.CENTER)
add_textbox(slide, "ANT-DBS\n術後12ヵ月\n平均発作減少率",
            Inches(0.4), Inches(2.75), Inches(3.5), Inches(0.95),
            font_size=14, color=WHITE, align=PP_ALIGN.CENTER)

# 棒グラフ的可視化（簡易）
timepoints = ["3ヵ月", "6ヵ月", "9ヵ月", "12ヵ月"]
dbs_vals   = [57.22,   61.61,   63.94,   65.28]
vns_vals   = [36.06,   39.94,   45.24,   48.35]

bar_area_x = Inches(4.3)
bar_area_y = Inches(1.4)
bar_w      = Inches(1.9)
bar_gap    = Inches(0.3)
max_h      = Inches(3.5)
scale      = max_h / 100

for i, tp in enumerate(timepoints):
    bx = bar_area_x + i * (bar_w * 2 + bar_gap + Inches(0.15))

    # VNS棒（薄）
    v_h = vns_vals[i] * float(scale)
    v_y = bar_area_y + max_h - v_h
    add_rect(slide, bx, v_y, bar_w * 0.85, v_h,
             RGBColor(0xBB, 0xCC, 0xDD))
    add_textbox(slide, f"{vns_vals[i]}%",
                bx, v_y - Inches(0.3), bar_w * 0.85, Inches(0.3),
                font_size=11, color=NAVY, align=PP_ALIGN.CENTER)

    # ANT-DBS棒（濃）
    d_h = dbs_vals[i] * float(scale)
    d_y = bar_area_y + max_h - d_h
    dx = bx + bar_w * 0.9
    add_rect(slide, dx, d_y, bar_w * 0.85, d_h, NAVY)
    add_textbox(slide, f"{dbs_vals[i]}%",
                dx, d_y - Inches(0.3), bar_w * 0.85, Inches(0.3),
                font_size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # 時点ラベル
    add_textbox(slide, f"術後\n{tp}",
                bx, bar_area_y + max_h + Inches(0.05),
                bar_w * 1.8, Inches(0.45),
                font_size=11, color=NAVY, align=PP_ALIGN.CENTER)

# 凡例
leg_y = Inches(5.15)
add_rect(slide, Inches(4.3), leg_y, Inches(0.35), Inches(0.25), RGBColor(0xBB, 0xCC, 0xDD))
add_textbox(slide, "VNS", Inches(4.7), leg_y, Inches(1.5), Inches(0.25), font_size=12, color=NAVY)
add_rect(slide, Inches(6.2), leg_y, Inches(0.35), Inches(0.25), NAVY)
add_textbox(slide, "ANT-DBS", Inches(6.6), leg_y, Inches(2.0), Inches(0.25), font_size=12, color=NAVY)

# 差分コメント
add_rect(slide, Inches(4.3), Inches(5.55), Inches(8.6), Inches(0.65), LIGHT_BG)
add_textbox(slide,
            "術後12ヵ月：ANT-DBS（65.28%）は VNS（48.35%）を 約17ポイント 上回る発作減少率を達成",
            Inches(4.4), Inches(5.62), Inches(8.4), Inches(0.5),
            font_size=14, bold=True, color=NAVY)

# ===========================================================
# スライド 5：レスポンダー率・発作消失率
# ===========================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
navy_header_bar(slide, "レスポンダー率・発作消失率（術後12ヵ月）")
elec_line(slide)
footer_bar(slide)
add_slide_number(slide, 5, TOTAL_SLIDES)

metrics = [
    ("レスポンダー率\n（発作50%以上減少）", "72.22%", "N=13/18", "58.82%", "N=10/17"),
    ("発作消失率",                           "22.22%", "N=4/18",  "17.65%", "N=3/17"),
    ("非レスポンダー率",                     "27.78%", "N=5/18",  "41.18%", "N=7/17"),
]

col_colors = [NAVY, RGBColor(0xBB, 0xCC, 0xDD)]
for mi, (label, dbs_v, dbs_n, vns_v, vns_n) in enumerate(metrics):
    row_y = Inches(1.5) + mi * Inches(1.7)

    # ラベル列
    add_rect(slide, Inches(0.4), row_y, Inches(3.6), Inches(1.5), LIGHT_BG)
    add_textbox(slide, label,
                Inches(0.5), row_y + Inches(0.3), Inches(3.4), Inches(0.9),
                font_size=14, bold=True, color=NAVY)

    # ANT-DBS 値
    add_rect(slide, Inches(4.2), row_y, Inches(4.0), Inches(1.5), NAVY)
    add_textbox(slide, dbs_v,
                Inches(4.2), row_y + Inches(0.05), Inches(4.0), Inches(0.9),
                font_size=38, bold=True, color=ELEC_BLUE if mi < 2 else WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, f"ANT-DBS  {dbs_n}",
                Inches(4.2), row_y + Inches(1.0), Inches(4.0), Inches(0.35),
                font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

    # VNS 値
    add_rect(slide, Inches(8.5), row_y, Inches(4.2), Inches(1.5), RGBColor(0xE0, 0xE8, 0xF0))
    add_textbox(slide, vns_v,
                Inches(8.5), row_y + Inches(0.05), Inches(4.2), Inches(0.9),
                font_size=38, bold=False, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(slide, f"VNS  {vns_n}",
                Inches(8.5), row_y + Inches(1.0), Inches(4.2), Inches(0.35),
                font_size=11, color=NAVY, align=PP_ALIGN.CENTER)

add_textbox(slide, "ANT-DBSはレスポンダー率・発作消失率ともにVNSを上回り、非レスポンダー率は低い",
            Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.45),
            font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ===========================================================
# スライド 6：サブグループ解析・患者選択
# ===========================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
navy_header_bar(slide, "サブグループ解析：ANT-DBSが特に有効な患者像")
elec_line(slide)
footer_bar(slide)
add_slide_number(slide, 6, TOTAL_SLIDES)

add_textbox(slide, "術前の臨床的特徴に基づいて、最適な治療を選択できる可能性があります",
            Inches(0.4), Inches(1.25), Inches(12.5), Inches(0.45),
            font_size=16, color=NAVY)

subgroups = [
    ("焦点性発作",     "71.15%",  "vs VNS 59.00%"),
    ("開頭術歴あり",   "81.20%",  "過去の開頭術後も高い効果"),
    ("罹病期間が長い", "有意な正の相関", "P < 0.05"),
]

for si, (sg_label, dbs_val, note) in enumerate(subgroups):
    bx = Inches(0.4) + si * Inches(4.25)
    by = Inches(2.0)
    add_rect(slide, bx, by, Inches(4.0), Inches(3.0), NAVY)
    add_textbox(slide, sg_label,
                bx + Inches(0.15), by + Inches(0.2), Inches(3.7), Inches(0.55),
                font_size=16, bold=True, color=ELEC_BLUE)
    add_rect(slide, bx + Inches(0.15), by + Inches(0.85), Inches(3.7), Inches(0.06), ELEC_BLUE)
    add_textbox(slide, dbs_val,
                bx, by + Inches(0.95), Inches(4.0), Inches(1.1),
                font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, note,
                bx + Inches(0.1), by + Inches(2.15), Inches(3.8), Inches(0.6),
                font_size=11, color=RGBColor(0xCC, 0xDD, 0xFF), italic=True)

# 患者選択ガイド
add_textbox(slide, "ANT-DBS 適応が考慮される患者の特徴：",
            Inches(0.4), Inches(5.25), Inches(12.5), Inches(0.4),
            font_size=15, bold=True, color=NAVY)

criteria = [
    "✔  焦点性発作",
    "✔  過去に開頭術（切除術）を受けている",
    "✔  罹病期間が長い（長期難治性）",
    "✔  手術時年齢が比較的高い（成人）",
]
for ci, c in enumerate(criteria):
    cx = Inches(0.5) + ci * Inches(3.1)
    add_rect(slide, cx, Inches(5.75), Inches(2.95), Inches(0.65),
             LIGHT_BG if ci % 2 == 0 else RGBColor(0xD8, 0xED, 0xF8))
    add_textbox(slide, c, cx + Inches(0.1), Inches(5.82), Inches(2.75), Inches(0.5),
                font_size=12, color=NAVY)

# ===========================================================
# スライド 7：安全性
# ===========================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
navy_header_bar(slide, "安全性：重篤な合併症なし、継続刺激で良好な忍容性")
elec_line(slide)
footer_bar(slide)
add_slide_number(slide, 7, TOTAL_SLIDES)

# 大コールアウト
add_rect(slide, Inches(0.4), Inches(1.35), Inches(12.4), Inches(1.0), NAVY)
add_textbox(slide, "観察期間中（術後12ヵ月）、両治療とも重篤な合併症・副作用は報告されませんでした",
            Inches(0.5), Inches(1.45), Inches(12.2), Inches(0.8),
            font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 安全性比較テーブル
headers = ["", "VNS", "ANT-DBS"]
rows = [
    ("主な有害事象", "嗄声、咳嗽、疼痛、呼吸困難", "植込み部疼痛、感染、リード偏位"),
    ("重篤な合併症", "なし", "なし"),
    ("刺激最適化",  "患者の忍容性に応じ段階的に調整", "患者の忍容性に応じ段階的に調整"),
]

col_x = [Inches(0.4), Inches(3.0), Inches(8.3)]
col_w = [Inches(2.55), Inches(5.2), Inches(4.9)]
header_y = Inches(2.55)

for hi, (hd, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    fill = NAVY if hi > 0 else LIGHT_BG
    add_rect(slide, cx, header_y, cw, Inches(0.5), fill)
    c = WHITE if hi > 0 else NAVY
    add_textbox(slide, hd, cx, header_y + Inches(0.06), cw, Inches(0.4),
                font_size=14, bold=True, color=c, align=PP_ALIGN.CENTER)

for ri, (row) in enumerate(rows):
    ry = header_y + Inches(0.55) + ri * Inches(0.8)
    fill_row = LIGHT_BG if ri % 2 == 0 else WHITE
    for ci, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
        add_rect(slide, cx, ry, cw, Inches(0.75),
                 fill_row if ci == 0 else (RGBColor(0xF5, 0xF5, 0xF5) if ri % 2 == 0 else WHITE))
        bold = (ci == 0)
        add_textbox(slide, cell, cx + Inches(0.1), ry + Inches(0.1), cw - Inches(0.2), Inches(0.55),
                    font_size=12, bold=bold, color=NAVY)

add_textbox(slide, "刺激パラメータは治療効果と患者の忍容性に応じて段階的に最適化",
            Inches(0.4), Inches(5.7), Inches(12.4), Inches(0.4),
            font_size=13, color=NAVY, italic=True)

# ===========================================================
# スライド 8：まとめ（訴求ポイント）
# ===========================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
navy_header_bar(slide, "まとめ：ANT-DBSが示す 3つの訴求ポイント")
elec_line(slide)
footer_bar(slide)
add_slide_number(slide, 8, TOTAL_SLIDES)

points = [
    ("1", "優れた有効性",
     "全観察時点でVNSを上回る発作減少率（術後12ヵ月：65.28% vs 48.35%）\n"
     "レスポンダー率72.22%・発作消失率22.22%を達成"),
    ("2", "患者に合わせた治療選択",
     "焦点性発作・開頭術歴あり・長期難治性 の患者で特に高い有効性\n"
     "術前データから治療選択の根拠が得られる"),
    ("3", "良好な忍容性と継続性",
     "重篤な合併症なく12ヵ月間継続可能\n"
     "効果は時間経過とともにさらに蓄積（3ヵ月→12ヵ月で改善）"),
]

for pi, (num, title, body) in enumerate(points):
    py = Inches(1.45) + pi * Inches(1.75)

    # 番号サークル風
    add_rect(slide, Inches(0.4), py, Inches(0.8), Inches(0.8), ELEC_BLUE)
    add_textbox(slide, num, Inches(0.4), py + Inches(0.05), Inches(0.8), Inches(0.7),
                font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(1.3), py, Inches(11.5), Inches(1.5),
             NAVY if pi % 2 == 0 else LIGHT_BG)
    tc = WHITE if pi % 2 == 0 else NAVY
    add_textbox(slide, title,
                Inches(1.5), py + Inches(0.1), Inches(11.1), Inches(0.45),
                font_size=17, bold=True, color=ELEC_BLUE if pi % 2 == 0 else NAVY)
    add_textbox(slide, body,
                Inches(1.5), py + Inches(0.55), Inches(11.1), Inches(0.8),
                font_size=13, color=tc)

# CTA
add_rect(slide, Inches(0.4), Inches(6.5), Inches(12.4), Inches(0.55), ELEC_BLUE)
add_textbox(slide,
            "薬剤抵抗性てんかんの患者さんに、次のステップを。  —  ANT-DBS 詳細はメドトロニック担当者まで",
            Inches(0.4), Inches(6.52), Inches(12.4), Inches(0.5),
            font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ===========================================================
# 保存
# ===========================================================
out_path = "/home/user/Shota/ANT-DBS_医局説明会スライド.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
